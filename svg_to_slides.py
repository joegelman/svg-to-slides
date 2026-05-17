#!/usr/bin/env python3
"""
svg_to_slides.py — Convert SVG files to PPTX with editable DrawingML vector paths.
Each SVG becomes a one-slide PPTX saved alongside the original file, ready to
import into Google Slides (File → Import slides).

Install deps: pip3 install --break-system-packages python-pptx lxml
Usage:        python3 ~/Documents/svg_to_slides.py file.svg [more.svg ...]
"""
import sys, re, math
from pathlib import Path
from xml.etree import ElementTree as ET

# Prefer locally-installed deps (put there by Install.command) over system packages
_lib = Path.home() / '.local' / 'share' / 'svg-to-slides' / 'lib'
if _lib.exists() and str(_lib) not in sys.path:
    sys.path.insert(0, str(_lib))

try:
    from pptx import Presentation
    from pptx.util import Emu
    from lxml import etree
except ImportError:
    sys.exit("Run Install.command, or: pip3 install python-pptx lxml")

# ── Constants ────────────────────────────────────────────────────────────────

SLIDE_W = 9144000   # 10 inches in EMU
SLIDE_H = 6858000   # 7.5 inches in EMU
COORD   = 100000    # DrawingML internal path coordinate space width

_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── Transform parsing ────────────────────────────────────────────────────────

def parse_transform(s):
    """Parse an SVG transform string into a single (x,y)→(x',y') function."""
    fns = []
    for m in re.finditer(r'(\w+)\(([^)]*)\)', s or ''):
        name = m.group(1)
        vals = [float(v) for v in re.split(r'[\s,]+', m.group(2).strip()) if v]
        if name == 'translate':
            tx = vals[0]; ty = vals[1] if len(vals) > 1 else 0.0
            fns.append(lambda x, y, tx=tx, ty=ty: (x+tx, y+ty))
        elif name == 'scale':
            sx = vals[0]; sy = vals[1] if len(vals) > 1 else sx
            fns.append(lambda x, y, sx=sx, sy=sy: (x*sx, y*sy))
        elif name == 'rotate':
            a = math.radians(vals[0])
            cx = vals[1] if len(vals) > 1 else 0.0
            cy = vals[2] if len(vals) > 2 else 0.0
            ca, sa = math.cos(a), math.sin(a)
            fns.append(lambda x, y, cx=cx, cy=cy, ca=ca, sa=sa:
                ((x-cx)*ca-(y-cy)*sa+cx, (x-cx)*sa+(y-cy)*ca+cy))
        elif name == 'matrix':
            a, b, c, d, e, f = vals[:6]
            fns.append(lambda x, y, a=a, b=b, c=c, d=d, e=e, f=f:
                (a*x+c*y+e, b*x+d*y+f))
    if not fns:
        return None
    # SVG "T1 T2" means combined matrix = T1·T2, so apply rightmost (T2) first
    def composed(x, y):
        for fn in reversed(fns):
            x, y = fn(x, y)
        return x, y
    return composed

def chain(outer, inner):
    """Compose two transform functions: apply inner first, then outer."""
    if outer is None: return inner
    if inner is None: return outer
    return lambda x, y: outer(*inner(x, y))

# ── SVG path tokenizer ───────────────────────────────────────────────────────

_CMD_ARGC = dict(M=2,m=2,L=2,l=2,H=1,h=1,V=1,v=1,
                 C=6,c=6,S=4,s=4,Q=4,q=4,T=2,t=2,A=7,a=7,Z=0,z=0)

def expand_path(d):
    """Tokenize & expand implicit repeats in an SVG path d attribute."""
    toks = re.findall(
        r'[MLHVCSQTAZmlhvcsqtaz]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?', d)
    segs, cmd, buf = [], None, []
    for t in toks:
        if t.isalpha():
            if cmd is not None: segs.append((cmd, buf))
            cmd, buf = t, []
        else:
            buf.append(float(t))
    if cmd is not None: segs.append((cmd, buf))

    out = []
    for cmd, args in segs:
        n = _CMD_ARGC.get(cmd.upper(), 0)
        if n == 0:
            out.append((cmd, [])); continue
        for i in range(0, max(len(args), n), n):
            chunk = args[i:i+n]
            if len(chunk) < n: break
            # After the first M/m, repeated coords become implicit L/l
            c = ('L' if cmd=='M' else 'l' if cmd=='m' else cmd) if i > 0 else cmd
            out.append((c, chunk))
    return out

# ── SVG tree traversal ───────────────────────────────────────────────────────

_NS_RE = re.compile(r'\{[^}]+\}')
def _tag(el): return _NS_RE.sub('', el.tag)

def collect(el, acc_xfm=None, inh_fill='#000000'):
    """Yield (d, transform_fn, fill_hex) for every <path> in the SVG tree."""
    local_xfm = parse_transform(el.get('transform', ''))
    xfm = chain(acc_xfm, local_xfm)

    style = el.get('style', '')
    m = re.search(r'fill:\s*([^;]+)', style)
    fill = (m.group(1).strip() if m else None) or el.get('fill') or inh_fill
    if fill in ('none', ''): fill = inh_fill

    if _tag(el) == 'path':
        d = el.get('d', '')
        if d:
            yield (d, xfm or (lambda x, y: (x, y)), fill)

    for child in el:
        yield from collect(child, xfm, fill)

# ── DrawingML path builder ───────────────────────────────────────────────────

def _el(name):    return etree.Element(f'{{{_A}}}{name}')
def _sub(p, name): return etree.SubElement(p, f'{{{_A}}}{name}')
def _pt(x, y):
    e = _el('pt'); e.set('x', str(x)); e.set('y', str(y)); return e

def make_a_path(d, xfm, vb_x, vb_y, vb_w, vb_h):
    """Convert an SVG path d string into an <a:path> lxml element."""
    cw = COORD
    ch = round(COORD * vb_h / vb_w) if vb_w else COORD

    def sc(x, y):
        tx, ty = xfm(x, y)
        return round((tx - vb_x) / vb_w * cw), round((ty - vb_y) / vb_h * ch)

    a_path = _el('path')
    a_path.set('w', str(cw))
    a_path.set('h', str(ch))

    cx = cy = mx = my = 0.0
    last_ctrl = None   # (x2,y2) from last bezier, for S/s reflection

    for cmd, args in expand_path(d):
        if cmd not in ('C','c','S','s','Q','q','T','t'):
            last_ctrl = None

        if cmd == 'M':
            cx, cy = mx, my = args
            _sub(a_path, 'moveTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'm':
            cx, cy = mx, my = cx+args[0], cy+args[1]
            _sub(a_path, 'moveTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'L':
            cx, cy = args
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'l':
            cx += args[0]; cy += args[1]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'H':
            cx = args[0]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'h':
            cx += args[0]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'V':
            cy = args[0]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'v':
            cy += args[0]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'C':
            x1,y1,x2,y2,x,y = args
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl = (x2, y2); cx, cy = x, y
        elif cmd == 'c':
            x1,y1,x2,y2,dx,dy = args
            ax1,ay1 = cx+x1, cy+y1
            ax2,ay2 = cx+x2, cy+y2
            ax, ay  = cx+dx, cy+dy
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(ax1,ay1))); e.append(_pt(*sc(ax2,ay2))); e.append(_pt(*sc(ax,ay)))
            last_ctrl = (ax2, ay2); cx, cy = ax, ay
        elif cmd == 'S':
            x2,y2,x,y = args
            lx,ly = last_ctrl if last_ctrl else (cx,cy)
            x1,y1 = 2*cx-lx, 2*cy-ly
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl = (x2, y2); cx, cy = x, y
        elif cmd == 's':
            dx2,dy2,dx,dy = args
            x2,y2 = cx+dx2, cy+dy2; x,y = cx+dx, cy+dy
            lx,ly = last_ctrl if last_ctrl else (cx,cy)
            x1,y1 = 2*cx-lx, 2*cy-ly
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl = (x2, y2); cx, cy = x, y
        elif cmd == 'Q':
            qx,qy,x,y = args
            # Convert quadratic to cubic bezier
            x1,y1 = cx + 2/3*(qx-cx), cy + 2/3*(qy-cy)
            x2,y2 = x  + 2/3*(qx-x),  y  + 2/3*(qy-y)
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl = (qx, qy); cx, cy = x, y
        elif cmd == 'q':
            dqx,dqy,dx,dy = args
            qx,qy = cx+dqx, cy+dqy; x,y = cx+dx, cy+dy
            x1,y1 = cx + 2/3*(qx-cx), cy + 2/3*(qy-cy)
            x2,y2 = x  + 2/3*(qx-x),  y  + 2/3*(qy-y)
            e = _sub(a_path, 'cubicBezTo')
            e.append(_pt(*sc(x1,y1))); e.append(_pt(*sc(x2,y2))); e.append(_pt(*sc(x,y)))
            last_ctrl = (qx, qy); cx, cy = x, y
        elif cmd in ('Z', 'z'):
            _sub(a_path, 'close')
            cx, cy = mx, my
        # A/a (arc): not supported — add a straight line to the endpoint as fallback
        elif cmd == 'A':
            cx, cy = args[5], args[6]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))
        elif cmd == 'a':
            cx += args[5]; cy += args[6]
            _sub(a_path, 'lnTo').append(_pt(*sc(cx, cy)))

    return a_path

# ── PPTX assembly ────────────────────────────────────────────────────────────

def hex_from_fill(fill_str):
    if fill_str in ('black', '#000000'): return '000000'
    if fill_str in ('white', '#ffffff'): return 'ffffff'
    if fill_str.startswith('#'):         return fill_str.lstrip('#').upper()
    return '000000'

def add_shapes(slide, paths, vb_x, vb_y, vb_w, vb_h):
    """One <p:sp> per SVG <path> — xfrm tightly fits each path's actual bounding box."""
    aspect = vb_w / vb_h if vb_h else 1
    if aspect >= SLIDE_W / SLIDE_H:
        sw = round(SLIDE_W * 0.85); sh = round(sw / aspect)
    else:
        sh = round(SLIDE_H * 0.85); sw = round(sh * aspect)
    ox = (SLIDE_W - sw) // 2
    oy = (SLIDE_H - sh) // 2

    # coordinate space used by make_a_path (must match)
    cw = COORD
    ch = round(COORD * vb_h / vb_w) if vb_w else COORD

    for sp_id, (d, xfm, fill_str) in enumerate(paths, start=2):
        hx = hex_from_fill(fill_str)
        a_path = make_a_path(d, xfm, vb_x, vb_y, vb_w, vb_h)

        # Compute tight bounding box from the generated <a:pt> elements,
        # shift all points to a (0,0) origin, then map to a tight EMU box.
        pts = list(a_path.iter(f'{{{_A}}}pt'))
        if pts:
            xs = [int(pt.get('x')) for pt in pts]
            ys = [int(pt.get('y')) for pt in pts]
            px_min, py_min = min(xs), min(ys)
            px_max, py_max = max(xs), max(ys)
            pw = max(px_max - px_min, 1)
            ph = max(py_max - py_min, 1)
            for pt in pts:
                pt.set('x', str(int(pt.get('x')) - px_min))
                pt.set('y', str(int(pt.get('y')) - py_min))
            a_path.set('w', str(pw))
            a_path.set('h', str(ph))
            shape_ox = ox + round(px_min / cw * sw)
            shape_oy = oy + round(py_min / ch * sh)
            shape_sw = max(round(pw / cw * sw), 1)
            shape_sh = max(round(ph / ch * sh), 1)
        else:
            shape_ox, shape_oy, shape_sw, shape_sh = ox, oy, sw, sh

        sp = etree.fromstring(
            f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{sp_id}" name="Path {sp_id - 1}"/>'
            f'<p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr>'
            f'<p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{shape_ox}" y="{shape_oy}"/><a:ext cx="{shape_sw}" cy="{shape_sh}"/></a:xfrm>'
            f'<a:custGeom>'
            f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
            f'<a:rect l="0" t="0" r="r" b="b"/>'
            f'<a:pathLst/>'
            f'</a:custGeom>'
            f'<a:solidFill><a:srgbClr val="{hx}"/></a:solidFill>'
            f'<a:ln><a:noFill/></a:ln>'
            f'</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
            f'</p:sp>'
        )
        pl = sp.find(f'.//{{{_A}}}pathLst')
        pl.append(a_path)
        slide.shapes._spTree.append(sp)

# ── Entry point ──────────────────────────────────────────────────────────────

def convert(svg_files):
    """Convert one or more SVG files into a single PPTX — one slide per SVG."""
    svg_files = [f for f in svg_files if Path(f).suffix.lower() == '.svg']
    if not svg_files:
        print('No .svg files provided.', file=sys.stderr)
        return

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for svg_file in svg_files:
        p = Path(svg_file)
        root = ET.parse(p).getroot()

        vb = root.get('viewBox', '0 0 100 100')
        vb_x, vb_y, vb_w, vb_h = [float(v) for v in re.split(r'[\s,]+', vb.strip())]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_shapes(slide, list(collect(root)), vb_x, vb_y, vb_w, vb_h)
        print(f'  + {p.name}')

    first = Path(svg_files[0])
    base = first.with_suffix('.pptx') if len(svg_files) == 1 else first.parent / 'slides.pptx'
    out = base
    n = 2
    while out.exists():
        out = base.parent / f'{base.stem} {n}{base.suffix}'
        n += 1
    prs.save(str(out))
    print(f'✓ {out}')

if __name__ == '__main__':
    if len(sys.argv) < 2:
        sys.exit(f'Usage: {Path(sys.argv[0]).name} file.svg [...]')
    try:
        convert(sys.argv[1:])
    except Exception as ex:
        print(f'✗ {ex}', file=sys.stderr)
        sys.exit(1)
